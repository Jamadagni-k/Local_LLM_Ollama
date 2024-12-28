import streamlit as st
import ollama
from PIL import Image
import tempfile
import pandas as pd
from docx import Document
from pptx import Presentation

st.title("💬 Personal AI asssistant ")

if "messages" not in st.session_state:
    st.session_state["messages"] = [{"role": "assistant", "content": "Hi! I'm a personal AI assistant. I'm here to help you with your queries offline. Jamie created me. What can I help you with today?"}]

### Write Message History
for msg in st.session_state.messages:
    if msg["role"] == "user":
        st.chat_message(msg["role"], avatar="🧑‍💻").write(msg["content"])
    else:
        st.chat_message(msg["role"], avatar="🤖").write(msg["content"])

## Generator for Streaming Tokens
def generate_response():
    response = ollama.chat(model='llama3.2', stream=True, messages=st.session_state.messages)
    for partial_resp in response:
        token = partial_resp["message"]["content"]
        st.session_state["full_message"] += token
        yield token

if prompt := st.chat_input():
    st.session_state.messages.append({"role": "user", "content": prompt})
    st.chat_message("user", avatar="🧑‍💻").write(prompt)
    st.session_state["full_message"] = ""
    st.chat_message("assistant", avatar="🤖").write_stream(generate_response)
    st.session_state.messages.append({"role": "assistant", "content": st.session_state["full_message"]})

# File uploader for text, images, videos, and Microsoft Office files
uploaded_file = st.file_uploader("Attach a file", type=["txt", "png", "jpg", "jpeg", "mp4", "avi", "mov", "xlsx", "xls", "docx", "pptx"])

if uploaded_file is not None:
    file_type = uploaded_file.type

    if "text" in file_type:
        text_content = uploaded_file.read().decode("utf-8")
        st.write("Uploaded Text File Content:")
        st.write(text_content)
        st.session_state.messages.append({"role": "user", "content": text_content})
    elif "image" in file_type:
        image = Image.open(uploaded_file)
        st.image(image, caption="Uploaded Image", use_column_width=True)
        st.session_state.messages.append({"role": "user", "content": "User uploaded an image."})
    elif "video" in file_type:
        tfile = tempfile.NamedTemporaryFile(delete=False) 
        tfile.write(uploaded_file.read())
        video_path = tfile.name
        st.video(video_path)
        st.session_state.messages.append({"role": "user", "content": "User uploaded a video."})
    elif "spreadsheet" in file_type:
        df = pd.read_excel(uploaded_file)
        st.write("Uploaded Excel File Content:")
        st.write(df)
        st.session_state.messages.append({"role": "user", "content": "User uploaded an Excel file."})
    elif "wordprocessingml" in file_type:
        doc = Document(uploaded_file)
        doc_text = "\n".join([para.text for para in doc.paragraphs])
        st.write("Uploaded Word Document Content:")
        st.write(doc_text)
        st.session_state.messages.append({"role": "user", "content": "User uploaded a Word document."})
    elif "presentation" in file_type:
        ppt = Presentation(uploaded_file)
        ppt_text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    ppt_text += shape.text + "\n"
        st.write("Uploaded PowerPoint Presentation Content:")
        st.write(ppt_text)
        st.session_state.messages.append({"role": "user", "content": "User uploaded a PowerPoint presentation."})
    else:
        st.warning("Unsupported file type.")
    
    # Generate and display the response for the uploaded file
    st.session_state["full_message"] = ""
    st.chat_message("assistant", avatar="🤖").write_stream(generate_response)
    st.session_state.messages.append({"role": "assistant", "content": st.session_state["full_message"]})
